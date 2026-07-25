#!/usr/bin/env python3
"""
PPTX文档转换器
将PPTX文档转换为HTML格式
"""

import logging
from typing import Tuple, Optional

from .base import BaseConverter, HTMLWrapper, ConversionError

logger = logging.getLogger(__name__)


class PptxConverter(BaseConverter):
    """PPTX文档转换器"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['pptx']
    
    def is_supported(self, file_type: str) -> bool:
        """检查文件类型是否支持转换"""
        return file_type.lower() in self.supported_formats
    
    def convert_to_html(self, file_path: str, file_type: str) -> Tuple[bool, Optional[str], Optional[str]]:
        """转换PPTX文档为HTML"""
        try:
            # 验证文件
            is_valid, error_msg = self.validate_file(file_path)
            if not is_valid:
                return False, None, error_msg

            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(file_path)
            content_parts = []

            # 添加演示文稿概览
            content_parts.append(f"<div class='presentation-overview'>")
            content_parts.append(f"<h1 class='presentation-title'>📽️ PowerPoint演示文稿</h1>")
            content_parts.append(f"<div class='presentation-info'>共 {len(prs.slides)} 张幻灯片</div>")
            content_parts.append(f"</div>")

            for i, slide in enumerate(prs.slides, 1):
                content_parts.append(f"<div class='slide-container' data-slide='{i}'>")
                content_parts.append(f"<div class='slide-header'>")
                content_parts.append(f"<h2 class='slide-title'>幻灯片 {i}</h2>")
                content_parts.append(f"<div class='slide-number'>第 {i} 页，共 {len(prs.slides)} 页</div>")
                content_parts.append(f"</div>")

                content_parts.append("<div class='slide-content'>")

                # 提取和处理各种类型的内容
                slide_content = self._extract_slide_content(slide)

                if slide_content['texts']:
                    # 处理文本内容
                    for j, text_info in enumerate(slide_content['texts']):
                        text = text_info['text']
                        is_title = text_info['is_title']

                        if is_title:
                            content_parts.append(f"<h3 class='slide-text-title'>{HTMLWrapper.escape_html(text)}</h3>")
                        else:
                            # 处理列表项
                            if text.startswith('•') or text.startswith('-') or text.startswith('*'):
                                content_parts.append(f"<ul class='slide-list'><li class='slide-list-item'>{HTMLWrapper.escape_html(text[1:].strip())}</li></ul>")
                            else:
                                content_parts.append(f"<p class='slide-text'>{HTMLWrapper.escape_html(text)}</p>")

                # 处理表格
                if slide_content['tables']:
                    for table_data in slide_content['tables']:
                        content_parts.append("<div class='slide-table-container'>")
                        content_parts.append("<table class='slide-table'>")
                        for row_data in table_data:
                            content_parts.append("<tr>")
                            for cell_text in row_data:
                                content_parts.append(f"<td class='slide-table-cell'>{HTMLWrapper.escape_html(cell_text)}</td>")
                            content_parts.append("</tr>")
                        content_parts.append("</table>")
                        content_parts.append("</div>")

                # 处理图片
                if slide_content['images']:
                    for img_info in slide_content['images']:
                        if img_info.get('placeholder', False):
                            # 显示占位符
                            content_parts.append(f"""
                            <div class='slide-image-placeholder'>
                                <div class='image-icon'>🖼️</div>
                                <div class='image-info'>图片: {img_info['name']}</div>
                                <div class='image-size'>{img_info['width']} × {img_info['height']}</div>
                            </div>
                            """)
                        else:
                            # 显示实际图片
                            content_parts.append(f"""
                            <div class='slide-image-container'>
                                <img src="data:image/{img_info['format']};base64,{img_info['base64']}"
                                     alt="{HTMLWrapper.escape_html(img_info['name'])}"
                                     class='slide-image'
                                     title="图片: {HTMLWrapper.escape_html(img_info['name'])} ({img_info['width']} × {img_info['height']})">
                                <div class='image-caption'>{HTMLWrapper.escape_html(img_info['name'])}</div>
                            </div>
                            """)

                # 如果没有任何内容
                if not any([slide_content['texts'], slide_content['tables'], slide_content['images']]):
                    content_parts.append("<p class='no-content'>此幻灯片无可显示内容</p>")

                content_parts.append("</div>")  # slide-content
                content_parts.append("</div>")  # slide-container

            html_content = "\n".join(content_parts)
            full_html = self._wrap_pptx_html(html_content, "PowerPoint文档预览")

            return True, full_html, None

        except Exception as e:
            logger.error(f"转换PPTX失败: {str(e)}")
            return False, None, str(e)

    def _extract_slide_content(self, slide):
        """提取幻灯片的各种内容"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        content = {
            'texts': [],
            'tables': [],
            'images': []
        }

        for shape in slide.shapes:
            try:
                # 处理文本框和标题
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # 判断是否为标题（基于位置和字体大小等）
                    is_title = self._is_title_shape(shape)
                    content['texts'].append({
                        'text': text,
                        'is_title': is_title
                    })

                # 处理表格
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_data = []
                    table = shape.table
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip() if cell.text else ""
                            row_data.append(cell_text)
                        table_data.append(row_data)
                    if table_data:
                        content['tables'].append(table_data)

                # 处理图片
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # 提取图片数据
                        image_data = self._extract_image_data(shape)
                        if image_data:
                            image_info = {
                                'name': getattr(shape, 'name', '未知图片'),
                                'width': f"{shape.width // 914400}cm" if hasattr(shape, 'width') else "未知",
                                'height': f"{shape.height // 914400}cm" if hasattr(shape, 'height') else "未知",
                                'data': image_data['data'],
                                'format': image_data['format'],
                                'base64': image_data['base64']
                            }
                            content['images'].append(image_info)
                        else:
                            # 如果无法提取图片数据，显示占位符
                            content['images'].append({
                                'name': getattr(shape, 'name', '未知图片'),
                                'width': f"{shape.width // 914400}cm" if hasattr(shape, 'width') else "未知",
                                'height': f"{shape.height // 914400}cm" if hasattr(shape, 'height') else "未知",
                                'placeholder': True
                            })
                    except Exception as e:
                        logger.debug(f"处理图片时出错: {str(e)}")
                        content['images'].append({
                            'name': '图片',
                            'width': '未知',
                            'height': '未知',
                            'placeholder': True
                        })

            except Exception as shape_error:
                logger.debug(f"处理形状时出错: {shape_error}")
                continue

        return content

    def _is_title_shape(self, shape):
        """判断形状是否为标题"""
        try:
            # 基于形状的位置判断（标题通常在上方）
            if hasattr(shape, 'top') and shape.top < 1000000:  # 约1cm
                return True

            # 基于占位符类型判断
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                placeholder_type = shape.placeholder_format.type
                # 标题和副标题的占位符类型
                if placeholder_type in [1, 2]:  # TITLE = 1, SUBTITLE = 2
                    return True

            # 基于文本长度判断（标题通常较短）
            if hasattr(shape, 'text') and len(shape.text.strip()) < 100:
                return True

        except:
            pass

        return False

    def _extract_image_data(self, shape):
        """提取图片数据"""
        try:
            import base64

            # 获取图片对象
            image = shape.image

            # 获取图片二进制数据
            image_bytes = image.blob

            # 确定图片格式
            image_format = 'png'  # 默认格式
            if image_bytes.startswith(b'\xff\xd8\xff'):
                image_format = 'jpeg'
            elif image_bytes.startswith(b'\x89PNG'):
                image_format = 'png'
            elif image_bytes.startswith(b'GIF'):
                image_format = 'gif'
            elif image_bytes.startswith(b'BM'):
                image_format = 'bmp'

            # 转换为base64
            base64_data = base64.b64encode(image_bytes).decode('utf-8')

            return {
                'data': image_bytes,
                'format': image_format,
                'base64': base64_data
            }

        except Exception as e:
            logger.debug(f"提取图片数据失败: {str(e)}")
            return None

    def _wrap_pptx_html(self, content: str, title: str) -> str:
        """为PPTX内容包装HTML页面"""
        escaped_title = HTMLWrapper.escape_html(title)
        
        return f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0, user-scalable=yes, minimum-scale=0.5, maximum-scale=3.0">
            <meta name="format-detection" content="telephone=no">
            <meta name="apple-mobile-web-app-capable" content="yes">
            <meta name="apple-mobile-web-app-status-bar-style" content="black-translucent">
            <title>{escaped_title}</title>
            <style>
                * {{
                    box-sizing: border-box;
                }}
                body {{
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
                    line-height: 1.6;
                    margin: 0;
                    padding: 5px;
                    background: #f5f5f5;
                    color: #333;
                    font-size: 14px;
                    -webkit-text-size-adjust: 100%;
                    -webkit-font-smoothing: antialiased;
                }}

                .presentation-overview {{
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    color: white;
                    padding: 20px;
                    margin-bottom: 20px;
                    border-radius: 8px;
                    text-align: center;
                    box-shadow: 0 4px 12px rgba(0,0,0,0.15);
                }}

                .presentation-title {{
                    margin: 0 0 10px 0;
                    font-size: 24px;
                    font-weight: 700;
                    text-shadow: 0 2px 4px rgba(0,0,0,0.3);
                }}

                .presentation-info {{
                    font-size: 16px;
                    opacity: 0.9;
                }}

                .slide-container {{
                    background: white;
                    margin-bottom: 20px;
                    border-radius: 8px;
                    box-shadow: 0 2px 8px rgba(0,0,0,0.1);
                    overflow: hidden;
                    page-break-after: always;
                }}

                .slide-header {{
                    background: linear-gradient(135deg, #ff6b35, #f7931e);
                    color: white;
                    padding: 15px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                }}

                .slide-title {{
                    margin: 0;
                    font-size: 18px;
                    font-weight: 600;
                    text-shadow: 0 1px 2px rgba(0,0,0,0.1);
                }}

                .slide-number {{
                    font-size: 12px;
                    opacity: 0.9;
                    background: rgba(255,255,255,0.2);
                    padding: 4px 8px;
                    border-radius: 12px;
                }}

                .slide-content {{
                    padding: 25px;
                    min-height: 200px;
                    background: white;
                }}

                .slide-text-title {{
                    color: #2c3e50;
                    font-size: 22px;
                    font-weight: 700;
                    margin: 0 0 20px 0;
                    padding-bottom: 10px;
                    border-bottom: 3px solid #3498db;
                    position: relative;
                }}

                .slide-text-title::after {{
                    content: '';
                    position: absolute;
                    bottom: -3px;
                    left: 0;
                    width: 50px;
                    height: 3px;
                    background: #e74c3c;
                }}

                .slide-text {{
                    margin: 15px 0;
                    line-height: 1.8;
                    font-size: 16px;
                    color: #444;
                    text-align: justify;
                }}

                .slide-list {{
                    margin: 15px 0;
                    padding-left: 0;
                    list-style: none;
                }}

                .slide-list-item {{
                    margin: 8px 0;
                    padding: 8px 0 8px 20px;
                    position: relative;
                    font-size: 15px;
                    line-height: 1.6;
                }}

                .slide-list-item::before {{
                    content: '▶';
                    position: absolute;
                    left: 0;
                    color: #3498db;
                    font-weight: bold;
                }}

                .slide-table-container {{
                    margin: 20px 0;
                    overflow-x: auto;
                    border-radius: 6px;
                    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                }}

                .slide-table {{
                    width: 100%;
                    border-collapse: collapse;
                    background: white;
                }}

                .slide-table-cell {{
                    padding: 12px 15px;
                    border: 1px solid #e0e0e0;
                    text-align: left;
                    vertical-align: top;
                    font-size: 14px;
                }}

                .slide-table tr:first-child .slide-table-cell {{
                    background: #f8f9fa;
                    font-weight: 600;
                    color: #495057;
                }}

                .slide-table tr:nth-child(even) .slide-table-cell {{
                    background: #fdfdfe;
                }}

                /* 图片容器样式 */
                .slide-image-container {{
                    margin: 20px 0;
                    text-align: center;
                    background: #f8f9fa;
                    border-radius: 8px;
                    padding: 15px;
                    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                }}

                .slide-image {{
                    max-width: 100%;
                    max-height: 400px;
                    height: auto;
                    border-radius: 4px;
                    box-shadow: 0 2px 8px rgba(0,0,0,0.15);
                }}

                .image-caption {{
                    margin-top: 10px;
                    font-size: 14px;
                    color: #666;
                    font-style: italic;
                }}

                /* 图片占位符样式 */
                .slide-image-placeholder {{
                    margin: 20px 0;
                    padding: 30px;
                    background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                    border-radius: 8px;
                    text-align: center;
                    color: white;
                    box-shadow: 0 4px 8px rgba(0,0,0,0.1);
                }}

                .image-icon {{
                    font-size: 48px;
                    margin-bottom: 10px;
                }}

                .image-info {{
                    font-size: 16px;
                    font-weight: 600;
                    margin-bottom: 5px;
                }}

                .image-size {{
                    font-size: 14px;
                    opacity: 0.9;
                }}

                .no-content {{
                    text-align: center;
                    color: #6c757d;
                    font-style: italic;
                    padding: 40px 20px;
                    background: #f8f9fa;
                    border-radius: 4px;
                    border: 2px dashed #dee2e6;
                }}
                
                /* 移动端优化 */
                @media (max-width: 768px) {{
                    body {{
                        padding: 3px;
                        font-size: 13px;
                    }}
                    
                    .slide-container {{
                        margin-bottom: 15px;
                        border-radius: 6px;
                    }}
                    
                    .slide-header {{
                        padding: 12px 15px;
                        flex-direction: column;
                        align-items: flex-start;
                        gap: 8px;
                    }}
                    
                    .slide-title {{
                        font-size: 16px;
                    }}
                    
                    .slide-number {{
                        font-size: 11px;
                        align-self: flex-end;
                    }}
                    
                    .slide-content {{
                        padding: 20px 15px;
                        min-height: 150px;
                    }}
                    
                    .slide-text-title {{
                        font-size: 18px;
                        margin-bottom: 12px;
                    }}
                    
                    .slide-text {{
                        font-size: 14px;
                        margin: 10px 0;
                    }}
                    
                    .no-content {{
                        padding: 30px 15px;
                        font-size: 13px;
                    }}
                }}
                
                /* 超小屏幕优化 */
                @media (max-width: 480px) {{
                    body {{
                        padding: 2px;
                        font-size: 12px;
                    }}
                    
                    .slide-header {{
                        padding: 10px 12px;
                    }}
                    
                    .slide-title {{
                        font-size: 15px;
                    }}
                    
                    .slide-content {{
                        padding: 15px 12px;
                    }}
                    
                    .slide-text-title {{
                        font-size: 16px;
                    }}
                    
                    .slide-text {{
                        font-size: 13px;
                    }}
                }}
                
                /* 打印样式 */
                @media print {{
                    body {{
                        background: white;
                        padding: 0;
                    }}
                    
                    .slide-container {{
                        box-shadow: none;
                        border: 1px solid #ddd;
                        page-break-after: always;
                    }}
                    
                    .slide-header {{
                        background: #f8f9fa !important;
                        color: #333 !important;
                    }}
                }}
            </style>
        </head>
        <body>
            {content}
            
            <script>
                document.addEventListener('DOMContentLoaded', function() {{
                    console.log('PowerPoint预览页面加载完成');
                    
                    // 添加幻灯片导航功能
                    const slides = document.querySelectorAll('.slide-container');
                    console.log(`共找到 ${{slides.length}} 张幻灯片`);
                    
                    // 为每个幻灯片添加点击事件（可选功能）
                    slides.forEach((slide, index) => {{
                        slide.addEventListener('click', function() {{
                            console.log(`点击了第 ${{index + 1}} 张幻灯片`);
                            // 这里可以添加更多交互功能
                        }});
                    }});
                    
                    // 键盘导航支持
                    document.addEventListener('keydown', function(e) {{
                        if (e.key === 'ArrowDown' || e.key === 'PageDown') {{
                            window.scrollBy(0, window.innerHeight * 0.8);
                            e.preventDefault();
                        }} else if (e.key === 'ArrowUp' || e.key === 'PageUp') {{
                            window.scrollBy(0, -window.innerHeight * 0.8);
                            e.preventDefault();
                        }}
                    }});
                }});
            </script>
        </body>
        </html>
        """
